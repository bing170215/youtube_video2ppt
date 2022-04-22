#@Author:   Casserole fish
#@Time:    2022/4/21 23:14

from pptx import Presentation  
from pptx.util import Inches
from pathlib import Path
import re
from PIL import Image
def natural_sort_key(s, _nsre=re.compile('([0-9]+)')):
    return [
        int(text)
        if text.isdigit() else text.lower()
        for text in _nsre.split(s)]

# ap = argparse.ArgumentParser()
# ap.add_argument("-p", "--path", required=False, default='', help="Path to load img files")
# args = vars(ap.parse_args())

def is_read_successfully(file):
    try:
        imgFile = Image.open(file)
        return True
    except Exception:
        return False

def create_ppt(dir_path, ppt_path):
    prs = Presentation()  
    blank_slide_layout = prs.slide_layouts[6]  

    dir_path = Path(dir_path)
    # dir_path = Path(args['path'])

    images = [img.name for img in dir_path.iterdir() if img.suffix==".jpg"]

    sorted_images = sorted(images, key=natural_sort_key)

    for img in sorted_images:
        if is_read_successfully(f"{dir_path/Path(img)}"):
            slide = prs.slides.add_slide(blank_slide_layout)
            pic = slide.shapes.add_picture(f"{dir_path/Path(img)}", Inches(0.5), Inches(0.75), width=Inches(9.2), height=Inches(6))

    # 删除掉临时文件
    for f in dir_path.glob("*"):
        if f.name == ppt_path.split('\\')[-1].split('.')[0]+'.mp4':
            f.unlink()
        if re.match('frame\d+.jpg',f.name):
            f.unlink()
    prs.save(ppt_path)

